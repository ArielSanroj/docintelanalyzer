"""
Sistema RAG Multimodal mejorado basado en RAG-Anything.
Implementa procesamiento de documentos multimodales con soporte para texto, imágenes, tablas y ecuaciones.
"""

import os
import logging
import json
import base64
from typing import List, Dict, Tuple, Optional, Any, Union
from dataclasses import dataclass, field
import numpy as np
from sentence_transformers import SentenceTransformer
from sklearn.metrics.pairwise import cosine_similarity
import re
from PIL import Image
import io
import requests
from pathlib import Path

# Importaciones para procesamiento de documentos
try:
    import fitz  # PyMuPDF
    PYMUPDF_AVAILABLE = True
except ImportError:
    PYMUPDF_AVAILABLE = False

try:
    import pandas as pd
    PANDAS_AVAILABLE = True
except ImportError:
    PANDAS_AVAILABLE = False

try:
    import cv2
    OPENCV_AVAILABLE = True
except ImportError:
    OPENCV_AVAILABLE = False

logger = logging.getLogger(__name__)

@dataclass
class MultimodalChunk:
    """Representa un fragmento multimodal del documento."""
    content: str
    content_type: str  # "text", "image", "table", "equation", "chart"
    start_pos: int
    end_pos: int
    chunk_id: int
    metadata: Dict = field(default_factory=dict)
    embedding: Optional[np.ndarray] = None
    visual_features: Optional[Dict] = None  # Para características visuales

@dataclass
class KnowledgeGraph:
    """Grafo de conocimiento para entidades y relaciones."""
    entities: Dict[str, Dict] = field(default_factory=dict)
    relations: List[Dict] = field(default_factory=list)
    entity_embeddings: Dict[str, np.ndarray] = field(default_factory=dict)

@dataclass
class MultimodalRetrievalResult:
    """Resultado de recuperación multimodal."""
    chunks: List[MultimodalChunk]
    scores: List[float]
    query: str
    total_chunks_searched: int
    retrieval_type: str
    cross_modal_matches: List[Dict] = field(default_factory=list)

class MultimodalRAGSystem:
    """Sistema RAG Multimodal basado en RAG-Anything."""
    
    def __init__(self, 
                 embedding_model: str = "all-MiniLM-L6-v2",
                 vlm_model: str = "nvidia/llama-3.2-vision-instruct",
                 enable_vision: bool = True):
        """
        Inicializa el sistema RAG multimodal.
        
        Args:
            embedding_model: Modelo de embeddings para texto
            vlm_model: Modelo de visión-lenguaje para análisis multimodal
            enable_vision: Habilitar capacidades de visión
        """
        self.embedding_model_name = embedding_model
        self.vlm_model_name = vlm_model
        self.enable_vision = enable_vision
        
        # Componentes del sistema
        self.embedding_model = None
        self.vlm_model = None
        self.knowledge_graph = KnowledgeGraph()
        self.multimodal_chunks: List[MultimodalChunk] = []
        self.chunk_embeddings: np.ndarray = None
        
        # Procesadores especializados
        self.text_processor = None
        self.image_processor = None
        self.table_processor = None
        self.equation_processor = None
        
        self._initialize_models()
        self._initialize_processors()
    
    def _initialize_models(self):
        """Inicializa los modelos necesarios."""
        try:
            logger.info("Inicializando modelos multimodales...")
            
            # Modelo de embeddings para texto
            cache_folder = os.path.expanduser("~/.cache/sentence_transformers")
            os.makedirs(cache_folder, exist_ok=True)
            
            self.embedding_model = SentenceTransformer(
                self.embedding_model_name,
                cache_folder=cache_folder
            )
            
            # Modelo VLM si está habilitado
            if self.enable_vision:
                try:
                    from langchain_nvidia_ai_endpoints import ChatNVIDIA
                    self.vlm_model = ChatNVIDIA(
                        model=self.vlm_model_name,
                        temperature=0.3,
                        max_tokens=4096
                    )
                    logger.info("Modelo VLM inicializado")
                except Exception as e:
                    logger.warning(f"No se pudo inicializar VLM: {e}")
                    self.enable_vision = False
            
            logger.info("Modelos multimodales inicializados correctamente")
            
        except Exception as e:
            logger.error(f"Error inicializando modelos: {e}")
            raise
    
    def _initialize_processors(self):
        """Inicializa los procesadores especializados."""
        self.text_processor = TextProcessor()
        self.image_processor = ImageProcessor() if self.enable_vision else None
        self.table_processor = TableProcessor()
        self.equation_processor = EquationProcessor()
    
    def process_document(self, file_path: str, content: Optional[str] = None) -> None:
        """
        Procesa un documento multimodal completo.
        
        Args:
            file_path: Ruta del archivo
            content: Contenido del documento (opcional)
        """
        logger.info(f"Procesando documento multimodal: {file_path}")
        
        # Detectar tipo de archivo
        file_ext = Path(file_path).suffix.lower()
        
        if file_ext == '.pdf':
            self._process_pdf(file_path)
        elif file_ext in ['.docx', '.doc']:
            self._process_word_document(file_path)
        elif file_ext in ['.pptx', '.ppt']:
            self._process_presentation(file_path)
        elif file_ext in ['.xlsx', '.xls']:
            self._process_spreadsheet(file_path)
        elif file_ext in ['.png', '.jpg', '.jpeg', '.gif', '.bmp']:
            self._process_image(file_path)
        elif content:
            self._process_text_content(content)
        else:
            raise ValueError(f"Formato de archivo no soportado: {file_ext}")
        
        # Construir grafo de conocimiento
        self._build_knowledge_graph()
        
        # Generar embeddings
        self._generate_multimodal_embeddings()
        
        logger.info(f"Documento procesado: {len(self.multimodal_chunks)} chunks multimodales")
    
    def _process_pdf(self, file_path: str):
        """Procesa un documento PDF con extracción multimodal."""
        if not PYMUPDF_AVAILABLE:
            raise ImportError("PyMuPDF no está disponible. Instale con: pip install PyMuPDF")
        
        logger.info("Procesando PDF con extracción multimodal...")
        
        doc = fitz.open(file_path)
        chunks = []
        
        for page_num in range(len(doc)):
            page = doc[page_num]
            
            # Extraer texto
            text = page.get_text()
            if text.strip():
                chunks.append(MultimodalChunk(
                    content=text,
                    content_type="text",
                    start_pos=len(chunks),
                    end_pos=len(chunks) + len(text),
                    chunk_id=len(chunks),
                    metadata={"page": page_num, "source": "pdf_text"}
                ))
            
            # Extraer imágenes
            if self.enable_vision:
                image_list = page.get_images()
                for img_index, img in enumerate(image_list):
                    try:
                        xref = img[0]
                        pix = fitz.Pixmap(doc, xref)
                        if pix.n - pix.alpha < 4:  # GRAY or RGB
                            img_data = pix.tobytes("png")
                            
                            # Procesar imagen
                            processed_img = self.image_processor.process_image(img_data)
                            
                            chunks.append(MultimodalChunk(
                                content=processed_img["description"],
                                content_type="image",
                                start_pos=len(chunks),
                                end_pos=len(chunks) + len(processed_img["description"]),
                                chunk_id=len(chunks),
                                metadata={
                                    "page": page_num,
                                    "image_index": img_index,
                                    "source": "pdf_image",
                                    "image_data": img_data
                                },
                                visual_features=processed_img
                            ))
                        pix = None
                    except Exception as e:
                        logger.warning(f"Error procesando imagen {img_index} en página {page_num}: {e}")
            
            # Extraer tablas
            tables = page.find_tables()
            for table_index, table in enumerate(tables):
                try:
                    table_data = table.extract()
                    if table_data:
                        processed_table = self.table_processor.process_table(table_data)
                        
                        chunks.append(MultimodalChunk(
                            content=processed_table["description"],
                            content_type="table",
                            start_pos=len(chunks),
                            end_pos=len(chunks) + len(processed_table["description"]),
                            chunk_id=len(chunks),
                            metadata={
                                "page": page_num,
                                "table_index": table_index,
                                "source": "pdf_table",
                                "table_data": table_data
                            }
                        ))
                except Exception as e:
                    logger.warning(f"Error procesando tabla {table_index} en página {page_num}: {e}")
        
        doc.close()
        self.multimodal_chunks = chunks
    
    def _process_word_document(self, file_path: str):
        """Procesa un documento de Word."""
        try:
            from docx import Document
            doc = Document(file_path)
            chunks = []
            
            for para in doc.paragraphs:
                if para.text.strip():
                    chunks.append(MultimodalChunk(
                        content=para.text,
                        content_type="text",
                        start_pos=len(chunks),
                        end_pos=len(chunks) + len(para.text),
                        chunk_id=len(chunks),
                        metadata={"source": "word_text"}
                    ))
            
            # Procesar tablas en Word
            for table in doc.tables:
                table_data = []
                for row in table.rows:
                    row_data = [cell.text for cell in row.cells]
                    table_data.append(row_data)
                
                if table_data:
                    processed_table = self.table_processor.process_table(table_data)
                    chunks.append(MultimodalChunk(
                        content=processed_table["description"],
                        content_type="table",
                        start_pos=len(chunks),
                        end_pos=len(chunks) + len(processed_table["description"]),
                        chunk_id=len(chunks),
                        metadata={"source": "word_table", "table_data": table_data}
                    ))
            
            self.multimodal_chunks = chunks
            
        except ImportError:
            logger.error("python-docx no está disponible. Instale con: pip install python-docx")
            raise
    
    def _process_presentation(self, file_path: str):
        """Procesa una presentación PowerPoint."""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            chunks = []
            
            for slide_num, slide in enumerate(prs.slides):
                # Procesar texto de la diapositiva
                slide_text = ""
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text += shape.text + "\n"
                
                if slide_text.strip():
                    chunks.append(MultimodalChunk(
                        content=slide_text,
                        content_type="text",
                        start_pos=len(chunks),
                        end_pos=len(chunks) + len(slide_text),
                        chunk_id=len(chunks),
                        metadata={"slide": slide_num, "source": "ppt_text"}
                    ))
            
            self.multimodal_chunks = chunks
            
        except ImportError:
            logger.error("python-pptx no está disponible. Instale con: pip install python-pptx")
            raise
    
    def _process_spreadsheet(self, file_path: str):
        """Procesa una hoja de cálculo Excel."""
        if not PANDAS_AVAILABLE:
            raise ImportError("pandas no está disponible. Instale con: pip install pandas")
        
        try:
            # Leer todas las hojas
            excel_file = pd.ExcelFile(file_path)
            chunks = []
            
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                
                # Procesar como tabla
                table_data = df.values.tolist()
                processed_table = self.table_processor.process_table(table_data)
                
                chunks.append(MultimodalChunk(
                    content=processed_table["description"],
                    content_type="table",
                    start_pos=len(chunks),
                    end_pos=len(chunks) + len(processed_table["description"]),
                    chunk_id=len(chunks),
                    metadata={"sheet": sheet_name, "source": "excel_table", "table_data": table_data}
                ))
            
            self.multimodal_chunks = chunks
            
        except Exception as e:
            logger.error(f"Error procesando Excel: {e}")
            raise
    
    def _process_image(self, file_path: str):
        """Procesa una imagen."""
        if not self.enable_vision:
            logger.warning("Procesamiento de imágenes deshabilitado")
            return
        
        try:
            with open(file_path, 'rb') as f:
                img_data = f.read()
            
            processed_img = self.image_processor.process_image(img_data)
            
            self.multimodal_chunks = [MultimodalChunk(
                content=processed_img["description"],
                content_type="image",
                start_pos=0,
                end_pos=len(processed_img["description"]),
                chunk_id=0,
                metadata={"source": "image_file", "file_path": file_path},
                visual_features=processed_img
            )]
            
        except Exception as e:
            logger.error(f"Error procesando imagen: {e}")
            raise
    
    def _process_text_content(self, content: str):
        """Procesa contenido de texto plano."""
        # Dividir en chunks de texto
        text_chunks = self.text_processor.split_text(content)
        
        self.multimodal_chunks = []
        for i, chunk_text in enumerate(text_chunks):
            self.multimodal_chunks.append(MultimodalChunk(
                content=chunk_text,
                content_type="text",
                start_pos=i * 1000,  # Aproximado
                end_pos=(i + 1) * 1000,
                chunk_id=i,
                metadata={"source": "text_content"}
            ))
    
    def _build_knowledge_graph(self):
        """Construye el grafo de conocimiento a partir de los chunks."""
        logger.info("Construyendo grafo de conocimiento...")
        
        self.knowledge_graph = KnowledgeGraph()
        
        for chunk in self.multimodal_chunks:
            # Extraer entidades del texto
            entities = self._extract_entities(chunk.content)
            
            for entity in entities:
                if entity not in self.knowledge_graph.entities:
                    self.knowledge_graph.entities[entity] = {
                        "type": "unknown",
                        "mentions": [],
                        "chunks": []
                    }
                
                self.knowledge_graph.entities[entity]["mentions"].append(chunk.chunk_id)
                self.knowledge_graph.entities[entity]["chunks"].append(chunk)
        
        # Extraer relaciones entre entidades
        self._extract_relations()
        
        logger.info(f"Grafo de conocimiento construido: {len(self.knowledge_graph.entities)} entidades, {len(self.knowledge_graph.relations)} relaciones")
    
    def _extract_entities(self, text: str) -> List[str]:
        """Extrae entidades del texto usando patrones simples."""
        entities = []
        
        # Patrones para diferentes tipos de entidades
        patterns = {
            "person": r'\b[A-Z][a-z]+ [A-Z][a-z]+\b',  # Nombres de personas
            "organization": r'\b[A-Z][a-z]+ (?:Inc|Corp|LLC|Ltd|Company|University|Institute)\b',
            "date": r'\b\d{1,2}[/-]\d{1,2}[/-]\d{2,4}\b',
            "number": r'\b\d+(?:\.\d+)?\b',
            "email": r'\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}\b',
            "url": r'https?://[^\s]+'
        }
        
        for entity_type, pattern in patterns.items():
            matches = re.findall(pattern, text)
            for match in matches:
                entities.append(match)
        
        return list(set(entities))  # Eliminar duplicados
    
    def _extract_relations(self):
        """Extrae relaciones entre entidades."""
        # Implementación simple de extracción de relaciones
        # En un sistema más avanzado, se usaría un modelo de NER/RE
        pass
    
    def _generate_multimodal_embeddings(self):
        """Genera embeddings para todos los chunks multimodales."""
        logger.info("Generando embeddings multimodales...")
        
        embeddings = []
        for chunk in self.multimodal_chunks:
            if chunk.content_type == "text":
                # Embedding de texto normal
                embedding = self.embedding_model.encode([chunk.content])[0]
            elif chunk.content_type == "image" and chunk.visual_features:
                # Para imágenes, usar descripción textual
                embedding = self.embedding_model.encode([chunk.content])[0]
            elif chunk.content_type == "table":
                # Para tablas, usar descripción textual
                embedding = self.embedding_model.encode([chunk.content])[0]
            else:
                # Fallback a embedding de texto
                embedding = self.embedding_model.encode([chunk.content])[0]
            
            embeddings.append(embedding)
            chunk.embedding = embedding
        
        self.chunk_embeddings = np.array(embeddings)
        logger.info(f"Embeddings multimodales generados: {self.chunk_embeddings.shape}")
    
    def retrieve_multimodal(self, query: str, top_k: int = 5) -> MultimodalRetrievalResult:
        """
        Recupera chunks relevantes usando búsqueda multimodal.
        
        Args:
            query: Consulta del usuario
            top_k: Número de chunks a recuperar
            
        Returns:
            Resultado de recuperación multimodal
        """
        logger.info(f"Recuperación multimodal para: '{query}'")
        
        # Generar embedding de la consulta
        query_embedding = self.embedding_model.encode([query])[0]
        
        # Calcular similitudes
        similarities = cosine_similarity([query_embedding], self.chunk_embeddings)[0]
        
        # Obtener top_k chunks
        top_indices = np.argsort(similarities)[::-1][:top_k]
        
        chunks = []
        scores = []
        cross_modal_matches = []
        
        for idx in top_indices:
            if similarities[idx] > 0.1:  # Umbral mínimo
                chunks.append(self.multimodal_chunks[idx])
                scores.append(float(similarities[idx]))
                
                # Detectar coincidencias cross-modal
                if self.multimodal_chunks[idx].content_type != "text":
                    cross_modal_matches.append({
                        "chunk_id": idx,
                        "content_type": self.multimodal_chunks[idx].content_type,
                        "relevance": float(similarities[idx])
                    })
        
        return MultimodalRetrievalResult(
            chunks=chunks,
            scores=scores,
            query=query,
            total_chunks_searched=len(self.multimodal_chunks),
            retrieval_type="multimodal",
            cross_modal_matches=cross_modal_matches
        )
    
    def get_document_stats(self) -> Dict:
        """Obtiene estadísticas del documento procesado."""
        if not self.multimodal_chunks:
            return {"status": "No document processed"}
        
        content_types = {}
        for chunk in self.multimodal_chunks:
            content_type = chunk.content_type
            content_types[content_type] = content_types.get(content_type, 0) + 1
        
        return {
            "total_chunks": len(self.multimodal_chunks),
            "content_types": content_types,
            "entities_count": len(self.knowledge_graph.entities),
            "relations_count": len(self.knowledge_graph.relations),
            "has_embeddings": self.chunk_embeddings is not None,
            "vision_enabled": self.enable_vision
        }


class TextProcessor:
    """Procesador especializado para texto."""
    
    def split_text(self, text: str, chunk_size: int = 1000, overlap: int = 200) -> List[str]:
        """Divide el texto en chunks con solapamiento."""
        if not text.strip():
            return []
        
        words = text.split()
        chunks = []
        
        start = 0
        while start < len(words):
            end = min(start + chunk_size, len(words))
            chunk = " ".join(words[start:end])
            chunks.append(chunk)
            start = end - overlap
        
        return chunks


class ImageProcessor:
    """Procesador especializado para imágenes."""
    
    def __init__(self):
        self.vlm_model = None
        self._initialize_vlm()
    
    def _initialize_vlm(self):
        """Inicializa el modelo VLM."""
        try:
            from langchain_nvidia_ai_endpoints import ChatNVIDIA
            self.vlm_model = ChatNVIDIA(
                model="nvidia/llama-3.2-vision-instruct",
                temperature=0.3,
                max_tokens=1024
            )
        except Exception as e:
            logger.warning(f"No se pudo inicializar VLM: {e}")
    
    def process_image(self, img_data: bytes) -> Dict:
        """Procesa una imagen y genera descripción."""
        try:
            # Convertir a base64 para el VLM
            img_base64 = base64.b64encode(img_data).decode('utf-8')
            
            if self.vlm_model:
                # Usar VLM para describir la imagen
                prompt = "Describe esta imagen en detalle, incluyendo texto visible, objetos, colores y composición."
                response = self.vlm_model.invoke([
                    {"type": "text", "text": prompt},
                    {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{img_base64}"}}
                ])
                description = response.content
            else:
                # Fallback a descripción básica
                description = "Imagen procesada (análisis detallado no disponible)"
            
            return {
                "description": description,
                "image_data": img_base64,
                "processed": True
            }
            
        except Exception as e:
            logger.error(f"Error procesando imagen: {e}")
            return {
                "description": "Error procesando imagen",
                "image_data": None,
                "processed": False
            }


class TableProcessor:
    """Procesador especializado para tablas."""
    
    def process_table(self, table_data: List[List[str]]) -> Dict:
        """Procesa una tabla y genera descripción."""
        try:
            if not table_data or not table_data[0]:
                return {"description": "Tabla vacía", "processed": False}
            
            # Generar descripción de la tabla
            rows = len(table_data)
            cols = len(table_data[0]) if table_data else 0
            
            # Extraer encabezados si existen
            headers = table_data[0] if table_data else []
            
            # Generar descripción
            description = f"Tabla con {rows} filas y {cols} columnas"
            if headers:
                description += f". Encabezados: {', '.join(headers[:5])}"  # Primeros 5 encabezados
                if len(headers) > 5:
                    description += "..."
            
            # Agregar algunas filas de ejemplo
            if len(table_data) > 1:
                sample_rows = table_data[1:3]  # Primeras 2 filas de datos
                description += f". Datos de ejemplo: {sample_rows}"
            
            return {
                "description": description,
                "table_data": table_data,
                "rows": rows,
                "cols": cols,
                "processed": True
            }
            
        except Exception as e:
            logger.error(f"Error procesando tabla: {e}")
            return {
                "description": "Error procesando tabla",
                "table_data": table_data,
                "processed": False
            }


class EquationProcessor:
    """Procesador especializado para ecuaciones matemáticas."""
    
    def process_equation(self, equation_text: str) -> Dict:
        """Procesa una ecuación matemática."""
        try:
            # Detectar si es LaTeX
            is_latex = any(marker in equation_text for marker in ['\\', '$$', '$'])
            
            description = f"Ecuación matemática"
            if is_latex:
                description += " en formato LaTeX"
            
            # Extraer variables y operadores
            variables = re.findall(r'[a-zA-Z][a-zA-Z0-9]*', equation_text)
            operators = re.findall(r'[+\-*/=<>]', equation_text)
            
            if variables:
                description += f" con variables: {', '.join(set(variables))}"
            if operators:
                description += f" y operadores: {', '.join(set(operators))}"
            
            return {
                "description": description,
                "equation_text": equation_text,
                "is_latex": is_latex,
                "variables": list(set(variables)),
                "operators": list(set(operators)),
                "processed": True
            }
            
        except Exception as e:
            logger.error(f"Error procesando ecuación: {e}")
            return {
                "description": "Error procesando ecuación",
                "equation_text": equation_text,
                "processed": False
            }